"""
PowerPoint exporter — generates an investment-grade deck.

Slides:
Equity mode (`AnalysisResult`)
  1. Title
  2. Company Overview
  3. Market & Competition
  4. Financial Snapshot
  5. Valuation Summary
  6. Investment Thesis
  7. Catalysts & Risks

M&A mode (`MAResult`)
  1. Title
  2. Opportunities / Pipeline
  3. Strategic Fit & Synergies
  4. Due Diligence Red Flags
  5. Deal Execution Plan
  6. LBO Snapshot (if relevant)
"""

from __future__ import annotations

from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..models import AcquisitionPipeline, AnalysisResult, FootballField, ICScoreSummary, MAResult, PeerCompsTable


# Simple, consistent styling (kept minimal to avoid template dependencies)
NAVY = RGBColor(27, 42, 74)
DARK = RGBColor(51, 51, 51)
MID = RGBColor(102, 102, 102)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(245, 247, 250)


def _set_run(run, *, bold: bool = False, size: int = 18, color=DARK) -> None:
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_header(slide, title: str, subtitle: str | None = None) -> None:
    # Light background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    # Top navy bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Title box
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.12), Inches(12.0), Inches(0.6))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run(r, bold=True, size=22, color=WHITE)
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.62), Inches(12.0), Inches(0.35))
        tf2 = sub_box.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        _set_run(r2, bold=False, size=12, color=MID)
        p2.alignment = PP_ALIGN.LEFT


def _add_bullets(slide, x, y, w, h, title: str, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    _set_run(r0, bold=True, size=20, color=DARK)

    for b in [b for b in bullets if b and str(b).strip()]:
        p = tf.add_paragraph()
        p.text = str(b).strip()
        p.level = 0
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"

def _add_table(slide, x, y, w, h, title: str, headers: list[str], rows: list[list[str]]) -> None:
    _add_bullets(slide, x, y, w, Inches(0.5), title, [])
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), x, y + Inches(0.5), w, h - Inches(0.5))
    table = table_shape.table

    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = head
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK
                p.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if i % 2 == 1 else RGBColor(250, 250, 250)


def _add_two_column(
    slide,
    *,
    left_title: str,
    left_bullets: list[str],
    right_title: str,
    right_bullets: list[str],
) -> None:
    _add_bullets(
        slide,
        Inches(0.75),
        Inches(1.8),
        Inches(6.0),
        Inches(5.2),
        left_title,
        left_bullets,
    )
    _add_bullets(
        slide,
        Inches(7.0),
        Inches(1.8),
        Inches(6.0),
        Inches(5.2),
        right_title,
        right_bullets,
    )


def _safe_list(xs: list[str] | None, fallback: str = "N/A") -> list[str]:
    if not xs:
        return [fallback]
    cleaned = [str(x).strip() for x in xs if x and str(x).strip()]
    return cleaned if cleaned else [fallback]


def _safe_str(s: str | None, fallback: str = "N/A") -> str:
    if s is None:
        return fallback
    s2 = str(s).strip()
    return s2 if s2 else fallback


def _build_equity_deck(result: AnalysisResult) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    f = result.fundamentals
    m = result.market
    fin = result.financials
    v = result.valuation
    t = result.thesis

    # 1) Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    subtitle = f"Generated {datetime.now().strftime('%B %d, %Y')} | { _safe_str(f.sector) }"
    _add_header(slide, _safe_str(f.company_name, result.company), subtitle)

    # Add key line
    key_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(12.2), Inches(1.5))
    tf = key_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r0 = p.add_run()
    r0.text = _safe_str(f.description, "Company description unavailable.")
    _set_run(r0, size=18, color=DARK)

    rec = (_safe_str(v.recommendation, "N/A")).upper()
    rec_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.6), Inches(12.2), Inches(0.6))
    tf2 = rec_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    _target = _safe_str(getattr(v, "target_price", None) or v.implied_value)
    _ev_label = f" | Implied EV: {_safe_str(v.implied_value)}" if getattr(v, "target_price", None) else ""
    r2.text = (
        f"Recommendation: {rec} | Target: {_target}{_ev_label} | "
        f"Upside/Downside: {_safe_str(v.upside_downside)} | Current: {_safe_str(v.current_price)}"
    )
    _set_run(r2, bold=True, size=14, color=NAVY)

    # 2) Company Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Company Overview")
    _add_two_column(
        slide,
        left_title="Business Model",
        left_bullets=[
            _safe_str(f.business_model, "N/A"),
            f"Founded: {_safe_str(f.founded)}",
            f"HQ: {_safe_str(f.headquarters)}",
            f"Employees: {_safe_str(f.employees)}",
            f"Ticker: {_safe_str(f.ticker)}",
        ],
        right_title="Competitive Advantages",
        right_bullets=_safe_list(f.competitive_advantages, "No advantages provided."),
    )

    # 3) Market & Competition
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Market & Competition")
    competitors = [f"{c.name} ({_safe_str(c.market_share, 'N/A')})" for c in (m.main_competitors or [])]
    _add_two_column(
        slide,
        left_title="Market Sizing",
        left_bullets=[
            f"TAM: {_safe_str(m.market_size)}",
            f"Growth (CAGR): {_safe_str(m.market_growth)}",
            f"Segment: {_safe_str(m.market_segment)}",
            f"Company share: {_safe_str(m.company_market_share)}",
        ],
        right_title="Trends & Competitors",
        right_bullets=_safe_list(m.key_trends, "No trends provided.") + (["Top competitors:"] + _safe_list(competitors, "N/A")),
    )

    # 4) Financial Snapshot (table)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Financial Snapshot")
    fin_rows = [
        ["Revenue", _safe_str(fin.revenue_current)],
        ["Revenue growth", _safe_str(fin.revenue_growth)],
        ["Gross margin", _safe_str(fin.gross_margin)],
        ["EBITDA margin", _safe_str(fin.ebitda_margin)],
        ["Net margin", _safe_str(fin.net_margin)],
        ["Free cash flow", _safe_str(fin.free_cash_flow)],
        ["Debt/Equity", _safe_str(fin.debt_to_equity)],
    ]
    _add_table(
        slide,
        Inches(0.75),
        Inches(1.4),
        Inches(6.2),
        Inches(5.7),
        "Key Financials",
        ["Metric", "Value"],
        fin_rows,
    )

    proj_rows = [
        [p.year, _safe_str(p.revenue), _safe_str(p.growth), _safe_str(p.ebitda_margin)]
        for p in (fin.projections or [])[:5]
    ] or [["—", "No projections provided.", "", ""]]
    _add_table(
        slide,
        Inches(7.1),
        Inches(1.4),
        Inches(5.45),
        Inches(5.7),
        "Outlook (Projections)",
        ["Year", "Revenue", "Growth", "EBITDA %"],
        proj_rows,
    )

    # 5) Valuation Summary (table)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Valuation Summary")
    methods_rows = [
        [
            meth.name,
            _safe_str(meth.low, "—"),
            _safe_str(meth.mid, "—"),
            _safe_str(meth.high, "—"),
            f"{meth.weight or '—'}%",
        ]
        for meth in (v.methods or [])
    ] or [["—", "—", "—", "—", "—"]]
    dcf = v.dcf_assumptions
    _tp = getattr(v, "target_price", None)
    _conclusion_bullets = [
        f"Implied EV: {_safe_str(v.implied_value)}",
    ]
    if _tp:
        _conclusion_bullets.append(f"Target price (per share): {_safe_str(_tp)}")
    _conclusion_bullets += [
        f"Upside/Downside: {_safe_str(v.upside_downside)}",
        f"Recommendation: {rec}",
        f"Current price: {_safe_str(v.current_price)}",
        f"DCF: WACC {_safe_str(getattr(dcf, 'wacc', None))}, TGR {_safe_str(getattr(dcf, 'terminal_growth', None))}",
    ]
    _add_bullets(
        slide,
        Inches(0.75),
        Inches(1.4),
        Inches(5.9),
        Inches(5.7),
        "Conclusion",
        _conclusion_bullets,
    )
    _add_table(
        slide,
        Inches(6.9),
        Inches(1.4),
        Inches(5.65),
        Inches(5.7),
        "Methods",
        ["Method", "Low", "Mid", "High", "Weight"],
        methods_rows,
    )

    # 6) Football Field — Bear / Base / Bull
    ff = result.football_field
    if ff:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_header(slide, "Valuation Football Field — Bear / Base / Bull")
        ff_rows = []
        for scenario in [ff.bear, ff.base, ff.bull]:
            if scenario:
                ff_rows.append([
                    scenario.name,
                    _safe_str(scenario.dcf_ev),
                    _safe_str(scenario.comps_ev),
                    _safe_str(scenario.blended_ev),
                    _safe_str(scenario.wacc),
                    _safe_str(scenario.ebitda_margin),
                ])
        if ff_rows:
            _add_table(
                slide,
                Inches(0.75), Inches(1.3), Inches(12.1), Inches(3.5),
                "Scenario Analysis",
                ["Scenario", "DCF EV", "Comps EV", "Blended EV", "WACC", "EBITDA Margin"],
                ff_rows,
            )
        _add_bullets(
            slide,
            Inches(0.75), Inches(5.0), Inches(12.2), Inches(2.2),
            "Ranges",
            [
                f"DCF range:     {_safe_str(ff.dcf_range)}",
                f"Comps range:   {_safe_str(ff.comps_range)}",
                f"Blended range: {_safe_str(ff.blended_range)}",
            ],
        )

    # 7) Peer Comparables
    pc = result.peer_comps
    if pc and pc.peers:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_header(
            slide,
            "Peer Comparables",
            f"{pc.n_peers} listed peers — real market data via yfinance",
        )
        peer_rows = [
            [
                p.name, p.ticker,
                _safe_str(p.ev_ebitda, "—"),
                _safe_str(p.ev_revenue, "—"),
                _safe_str(p.ebitda_margin, "—"),
                _safe_str(p.revenue_growth, "—"),
            ]
            for p in pc.peers
        ]
        peer_rows.append([
            "MEDIAN", "—",
            _safe_str(pc.median_ev_ebitda, "—"),
            _safe_str(pc.median_ev_revenue, "—"),
            _safe_str(pc.median_ebitda_margin, "—"),
            "—",
        ])
        _add_table(
            slide,
            Inches(0.75), Inches(1.3), Inches(12.1), Inches(5.8),
            "Trading Comps",
            ["Company", "Ticker", "EV/EBITDA", "EV/Revenue", "EBITDA Margin", "Rev Growth"],
            peer_rows,
        )

    # 8) IC Score Breakdown
    ic = result.ic_score
    if ic:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        rec_label = _safe_str(ic.recommendation)
        _add_header(slide, f"IC Score — {rec_label}", f"Investment Committee Scorecard")
        _add_table(
            slide,
            Inches(0.75), Inches(1.2), Inches(5.8), Inches(4.5),
            "Dimension Scores",
            ["Dimension", "Score"],
            [
                ["Strategy",    _safe_str(ic.strategy)],
                ["Synergies",   _safe_str(ic.synergies)],
                ["Financial",   _safe_str(ic.financial)],
                ["LBO",         _safe_str(ic.lbo)],
                ["Integration", _safe_str(ic.integration)],
                ["Risk",        _safe_str(ic.risk)],
                ["TOTAL",       _safe_str(ic.ic_score)],
            ],
        )
        _add_bullets(
            slide,
            Inches(7.0), Inches(1.2), Inches(5.8), Inches(4.5),
            "Rationale & Next Steps",
            [_safe_str(ic.rationale, "N/A")] + _safe_list(ic.next_steps),
        )

    # 9) Investment Thesis
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Investment Thesis")
    _add_bullets(
        slide,
        Inches(0.75),
        Inches(1.8),
        Inches(12.2),
        Inches(5.5),
        "Thesis",
        [
            _safe_str(t.thesis, "Thesis unavailable."),
            f"Bull: {_safe_str(t.bull_case)}",
            f"Base: {_safe_str(t.base_case)}",
            f"Bear: {_safe_str(t.bear_case)}",
        ],
    )

    # 10) Catalysts & Risks
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Catalysts & Risks")
    risks = [f"{r.level.upper()}: {r.text}" for r in (f.key_risks or [])]
    _add_two_column(
        slide,
        left_title="Catalysts",
        left_bullets=_safe_list(t.catalysts, "No catalysts provided."),
        right_title="Key Risks",
        right_bullets=_safe_list(risks, "No risks provided."),
    )

    return prs


def _build_ma_deck(result: MAResult) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    subtitle_parts = [
        f"Generated {datetime.now().strftime('%B %d, %Y')}",
        f"Target: {result.company} ({result.company_type})",
    ]
    if result.acquirer:
        subtitle_parts.append(f"Acquirer: {result.acquirer}")
    subtitle = " | ".join(subtitle_parts)

    # 1) Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "M&A Opportunity Deck", subtitle)

    obj = _safe_str(result.deal_sourcing.acquirer_objective, "Objective: N/A")
    headline_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(12.2), Inches(0.9))
    tf = headline_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r0 = p.add_run()
    r0.text = obj
    _set_run(r0, size=18, color=DARK, bold=True)

    stats = [
        f"Opportunities: {len(result.deal_sourcing.opportunities)}",
        f"Fit score: {_safe_str(result.strategic_fit.fit_score)}",
        f"Red flags: {len(result.due_diligence.red_flags)}",
        f"LBO feasible: {str(result.lbo.feasible) if result.lbo.feasible is not None else 'N/A'}",
    ]
    _add_bullets(slide, Inches(0.75), Inches(2.3), Inches(12.2), Inches(4.7), "Key Takeaways", stats)

    # 2) Opportunities / Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Opportunities / Pipeline")
    opp_rows = []
    for o in (result.deal_sourcing.opportunities or [])[:12]:
        rationale = _safe_str(o.rationale, "")
        rationale = (rationale[:140] + "…") if len(rationale) > 140 else rationale
        opp_rows.append(
            [
                _safe_str(o.name),
                _safe_str(o.geography),
                _safe_str(o.est_size),
                rationale or "N/A",
            ]
        )
    if not opp_rows:
        opp_rows = [["N/A", "N/A", "N/A", "No opportunities returned."]]
    _add_table(
        slide,
        Inches(0.75),
        Inches(1.2),
        Inches(12.1),
        Inches(6.3),
        "Pipeline",
        ["Name", "Geo", "Size", "Rationale (short)"],
        opp_rows,
    )

    # 3) Strategic Fit & Synergies
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Strategic Fit & Synergies")
    synergies = [
        f"{_safe_str(s.type).title()}: {_safe_str(s.description)} "
        f"(Impact: {_safe_str(s.est_impact)}, Timing: {_safe_str(s.timing)})"
        for s in (result.strategic_fit.key_synergies or [])[:6]
    ]
    _add_two_column(
        slide,
        left_title="Fit Summary",
        left_bullets=[
            f"Fit score: {_safe_str(result.strategic_fit.fit_score)}",
            f"Recommended structure: {_safe_str(result.strategic_fit.recommended_structure)}",
            f"Integration complexity: {_safe_str(result.strategic_fit.integration_complexity)}",
        ],
        right_title="Key Synergies",
        right_bullets=_safe_list(synergies, "No synergies provided."),
    )

    # 4) Due Diligence — Red Flags
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Due Diligence — Red Flags")
    rf_rows = []
    for rf in (result.due_diligence.red_flags or [])[:12]:
        finding = _safe_str(rf.finding, "")
        finding = (finding[:140] + "…") if len(finding) > 140 else finding
        rf_rows.append([_safe_str(rf.area), _safe_str(rf.severity).upper(), finding, _safe_str(rf.mitigation)])
    if not rf_rows:
        rf_rows = [["N/A", "N/A", "No red flags returned.", ""]]
    _add_table(
        slide,
        Inches(0.75),
        Inches(1.2),
        Inches(12.1),
        Inches(6.3),
        "Red Flags",
        ["Area", "Severity", "Finding (short)", "Mitigation"],
        rf_rows,
    )

    # 5) Deal Execution Plan
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Deal Execution Plan")
    _add_two_column(
        slide,
        left_title="Workplan",
        left_bullets=_safe_list(result.deal_execution.workplan, "No workplan provided."),
        right_title="Negotiation / Approvals",
        right_bullets=_safe_list(result.deal_execution.negotiation_points, "N/A")
        + ["Approvals:"]
        + _safe_list(result.deal_execution.approvals, "N/A"),
    )

    # 6) LBO Snapshot
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "LBO Snapshot (high-level)")
    _add_two_column(
        slide,
        left_title="Assumptions",
        left_bullets=[
            f"Feasible: {str(result.lbo.feasible) if result.lbo.feasible is not None else 'N/A'}",
            f"Entry multiple: {_safe_str(result.lbo.entry_multiple)}",
            f"Leverage: {_safe_str(result.lbo.leverage)}",
            f"Exit multiple: {_safe_str(result.lbo.exit_multiple)}",
            f"IRR range: {_safe_str(result.lbo.irr_range)}",
        ],
        right_title="Sensitivities",
        right_bullets=_safe_list(result.lbo.key_sensitivities, "N/A"),
    )

    return prs


def _build_pipeline_deck(result: AcquisitionPipeline) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    subtitle = f"Generated {datetime.now().strftime('%B %d, %Y')} | Buyer: {result.buyer}"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Acquisition Pipeline", subtitle)
    _add_bullets(
        slide,
        Inches(0.75),
        Inches(1.25),
        Inches(12.2),
        Inches(1.8),
        "Thesis / Focus",
        [_safe_str(result.thesis, "Thesis unavailable."), _safe_str(result.focus, "")],
    )
    _add_bullets(
        slide,
        Inches(0.75),
        Inches(3.2),
        Inches(12.2),
        Inches(4.0),
        "Screening Criteria",
        result.screening_criteria or ["N/A"],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Target Shortlist (with valuation estimates)")
    rows = []
    for t in (result.targets or [])[:14]:
        rows.append(
            [
                _safe_str(t.name),
                _safe_str(t.headquarters),
                _safe_str(t.segment),
                _safe_str(t.revenue_working),
                _safe_str(t.ebitda_margin),
                _safe_str(t.implied_ev),
            ]
        )
    if not rows:
        rows = [["N/A", "N/A", "N/A", "N/A", "N/A", "N/A"]]
    _add_table(
        slide,
        Inches(0.75),
        Inches(1.2),
        Inches(12.1),
        Inches(6.3),
        "Shortlist",
        ["Target", "HQ", "Segment", "Revenue (wk.)", "EBITDA %", "Implied EV"],
        rows,
    )

    # Add 1 slide per top targets (deep dive)
    for t in (result.targets or [])[:5]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_header(slide, _safe_str(t.name, "Target"))
        _add_two_column(
            slide,
            left_title="Investment Case",
            left_bullets=_safe_list(t.why_attractive, "N/A")
            + ["Strategic value:"] + _safe_list(t.strategic_value, "N/A"),
            right_title="Synergies / Risks",
            right_bullets=["Synergies:"] + _safe_list(t.synergies, "N/A")
            + ["Risks:"] + _safe_list(t.key_risks, "N/A"),
        )
        _add_bullets(
            slide,
            Inches(0.75),
            Inches(6.05),
            Inches(12.2),
            Inches(1.1),
            "Valuation (est.)",
            [
                f"Revenue range: {_safe_str(t.revenue_range)} | Working: {_safe_str(t.revenue_working)} | EBITDA: {_safe_str(t.ebitda_margin)} | EV: {_safe_str(t.implied_ev)}",
                "Rationale: " + "; ".join((t.valuation_rationale or [])[:3]),
            ],
        )

    return prs


def generate_pptx(result: AnalysisResult | MAResult | AcquisitionPipeline, output_path: str = "analysis.pptx") -> str:
    prs: Presentation
    if isinstance(result, AcquisitionPipeline):
        prs = _build_pipeline_deck(result)
    elif isinstance(result, MAResult):
        prs = _build_ma_deck(result)
    else:
        prs = _build_equity_deck(result)

    prs.save(output_path)
    return output_path
